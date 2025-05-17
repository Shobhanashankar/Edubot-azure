import os
import re
import requests
import fitz  # PyMuPDF
import nltk
import language_tool_python
from pptx import Presentation
from transformers import pipeline
from nltk.tokenize.punkt import PunktSentenceTokenizer, PunktParameters
from dotenv import load_dotenv
nltk.download('punkt')
load_dotenv()

subscription_key = os.getenv("VISION_KEY")
endpoint = os.getenv("VISION_ENDPOINT")
ocr_url = endpoint + "/vision/v3.2/ocr?language=en"

headers = {
    'Ocp-Apim-Subscription-Key': subscription_key,
    'Content-Type': 'application/octet-stream'
}

def correct_grammar(text):
    tool = None
    try:
        tool = language_tool_python.LanguageTool('en-US')
        matches = tool.check(text)
        return language_tool_python.utils.correct(text, matches)
    finally:
        if tool:
            tool.close()

def clean_text(raw_text):
    cleaned = re.sub(r'\s+', ' ', raw_text)
    cleaned = cleaned.replace('|', 'I').replace('0', 'O')
    cleaned = re.sub(r'[^\x00-\x7F]+', '', cleaned)
    cleaned = re.sub(r'[^\w\s.,!?-]', '', cleaned)

    lines = raw_text.splitlines()
    paragraph = ""
    paragraphs = []

    for line in lines:
        line = line.strip()
        if not line:
            if paragraph:
                paragraphs.append(paragraph.strip())
                paragraph = ""
        elif len(line) < 60:
            paragraph += " " + line
        else:
            if paragraph:
                paragraphs.append(paragraph.strip())
                paragraph = ""
            paragraphs.append(line)

    if paragraph:
        paragraphs.append(paragraph.strip())

    return "\n\n".join(paragraphs)

def chunk_text(text, max_chunk_size=1000):
    punkt_param = PunktParameters()
    tokenizer = PunktSentenceTokenizer(punkt_param)
    sentences = tokenizer.tokenize(text)

    chunks = []
    current_chunk = ""

    for sent in sentences:
        if len(current_chunk) + len(sent) + 1 <= max_chunk_size:
            current_chunk += " " + sent
        else:
            chunks.append(current_chunk.strip())
            current_chunk = sent
    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks

def generate_summary(text):
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    chunks = chunk_text(text)
    summaries = []
    for chunk in chunks:
        max_len = min(150, max(30, len(chunk) // 5))
        min_len = min(40, max_len // 2)
        result = summarizer(chunk, max_length=max_len, min_length=min_len, do_sample=False)
        summaries.append(result[0]['summary_text'])
    return " ".join(summaries).strip()

def ocr_image_bytes(image_bytes):
    response = requests.post(ocr_url, headers=headers, data=image_bytes)
    response.raise_for_status()
    return response.json()

def extract_text_from_ocr_result(result):
    text = ""
    for region in result.get("regions", []):
        for line in region.get("lines", []):
            line_text = " ".join([word["text"] for word in line["words"]])
            text += line_text + "\n"
    return text

def process_pdf_bytes(pdf_path):
    full_text = ""
    # Use context manager to ensure file is closed before deletion
    with fitz.open(pdf_path) as doc:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            img_bytes = pix.tobytes()
            ocr_result = ocr_image_bytes(img_bytes)
            page_text = extract_text_from_ocr_result(ocr_result)
            full_text += page_text + "\n"

    cleaned = clean_text(full_text)
    corrected = correct_grammar(cleaned)
    summary = generate_summary(corrected)

    return {
        "raw_text": full_text,
        "cleaned_text": cleaned,
        "corrected_text": corrected,
        "summary": summary
    }

def process_image_bytes(image_path):
    with open(image_path, 'rb') as image_file:
        image_bytes = image_file.read()

    result = ocr_image_bytes(image_bytes)
    text = extract_text_from_ocr_result(result)

    cleaned = clean_text(text)
    corrected = correct_grammar(cleaned)
    summary = generate_summary(corrected)

    return {
        "raw_text": text,
        "cleaned_text": cleaned,
        "corrected_text": corrected,
        "summary": summary
    }

def process_pptx_text(pptx_path):
    prs = Presentation(pptx_path)
    full_text = ""

    for i, slide in enumerate(prs.slides):
        slide_text = f"Slide {i + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        full_text += slide_text

    cleaned = clean_text(full_text)
    corrected = correct_grammar(cleaned)
    summary = generate_summary(corrected)

    return {
        "raw_text": full_text,
        "cleaned_text": cleaned,
        "corrected_text": corrected,
        "summary": summary
    }
